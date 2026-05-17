"""Generate the CSCI461 Healthcare Big Data project presentation."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ──────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0D, 0x3B, 0x6E)   # slide header bg
MID_BLUE    = RGBColor(0x15, 0x57, 0xA8)   # accent / table header
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF7)   # table alt-row
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT_GOLD = RGBColor(0xF5, 0xA6, 0x23)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

BLANK = prs.slide_layouts[6]          # completely blank layout


# ── helpers ──────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_rgb=None, line_rgb=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_TEXT,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def slide_header(slide, title_text, subtitle_text=""):
    """Dark-blue header bar at the top of content slides."""
    add_rect(slide, 0, 0, 13.33, 1.15, fill_rgb=DARK_BLUE)
    add_text(slide, title_text, 0.35, 0.12, 12.3, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle_text:
        add_text(slide, subtitle_text, 0.35, 0.78, 12.3, 0.35,
                 font_size=14, color=ACCENT_GOLD, align=PP_ALIGN.LEFT)
    # thin gold accent line under the header
    add_rect(slide, 0, 1.15, 13.33, 0.04, fill_rgb=ACCENT_GOLD)


def add_bullet_box(slide, items, l, t, w, h,
                   font_size=16, color=DARK_TEXT, title=None, title_color=MID_BLUE):
    """Add a box with optional bold title then bullet items."""
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True

    first = True
    if title:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = title
        run.font.bold = True
        run.font.size = Pt(font_size + 2)
        run.font.color.rgb = title_color
        first = False

    for item in items:
        p = tf.paragraphs[0] if (first and not title) else tf.add_paragraph()
        p.level = 0
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = ("• " if not item.startswith("  ") else "    ◦ ") + item.lstrip()
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        first = False
    return txb


def footer(slide, slide_num, total=10):
    add_rect(slide, 0, 7.15, 13.33, 0.35, fill_rgb=DARK_BLUE)
    add_text(slide, "CSCI461 — Introduction to Big Data  |  Spring 2026  |  Nile University",
             0.3, 7.17, 11, 0.30, font_size=11, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, f"{slide_num} / {total}",
             12.3, 7.17, 0.9, 0.30, font_size=11, color=WHITE, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)

# full background
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=DARK_BLUE)

# big gold accent stripe
add_rect(s, 0, 2.7, 13.33, 0.08, fill_rgb=ACCENT_GOLD)
add_rect(s, 0, 4.55, 13.33, 0.08, fill_rgb=ACCENT_GOLD)

# main title
add_text(s, "Real-Time Healthcare Adverse Event Analytics",
         0.6, 1.1, 12.1, 1.0,
         font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "using Big Data Streaming and Machine Learning",
         0.6, 2.05, 12.1, 0.7,
         font_size=26, bold=False, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# sub-info
add_text(s, "Youssif Zaghloul  •  Kazzy",
         0.6, 2.9, 12.1, 0.55,
         font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "CSCI461 — Introduction to Big Data   |   Spring 2026   |   Nile University",
         0.6, 3.5, 12.1, 0.45,
         font_size=16, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# tag-line
add_text(s, "Apache Kafka  ·  HDFS  ·  Spark MLlib  ·  PostgreSQL  ·  Metabase",
         0.6, 4.75, 12.1, 0.5,
         font_size=15, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

add_text(s, "1 / 10", 12.1, 7.1, 1.1, 0.35,
         font_size=11, color=WHITE, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT & MOTIVATION
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Problem Statement & Motivation",
             "Why does real-time pharmacovigilance matter?")

# two columns
add_rect(s, 0.35, 1.35, 6.0, 5.5, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_rect(s, 6.95, 1.35, 6.0, 5.5, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)

add_text(s, "The Challenge", 0.55, 1.4, 5.6, 0.45,
         font_size=17, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "FDA receives millions of spontaneous adverse-event reports each year via MedWatch",
    "Traditional pipelines are batch-oriented — reports analysed offline with significant latency",
    "Delayed detection of drug safety signals can cost lives",
    "No open-source, self-contained, end-to-end reference pipeline exists for course-level learning",
], 0.55, 1.85, 5.7, 4.8, font_size=14)

add_text(s, "Our Solution", 7.15, 1.4, 5.6, 0.45,
         font_size=17, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Ingest openFDA Drug Adverse Event records in near real-time via Apache Kafka",
    "Store raw + cleaned data in HDFS (bronze-silver-gold architecture)",
    "Apply Spark MLlib risk classification to 7,573+ records",
    "Surface insights via PostgreSQL + Metabase dashboards",
    "Fully containerised — runs with a single docker compose up",
], 7.15, 1.85, 5.7, 4.8, font_size=14)

footer(s, 2)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — PROJECT OBJECTIVES
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Project Objectives", "SMART goals aligned with the CSCI461 rubric")

# 4 objective cards
card_data = [
    ("Ingestion",     "Stream >7,500 FDA adverse-event records through Kafka within a single ingest session using the openFDA REST API."),
    ("Storage",       "Implement a bronze-silver-gold HDFS data lakehouse with automated Spark Structured Streaming cleaning and deduplication."),
    ("Classification","Train and evaluate Logistic Regression and Random Forest models (+ rule-based baseline) achieving ≥0.85 weighted F1 on risk labels."),
    ("Reporting",     "Aggregate and visualise risk predictions by patient sex and medicinal product in a self-service Metabase dashboard."),
]

xs = [0.35, 6.84, 0.35, 6.84]
ys = [1.35, 1.35, 4.25, 4.25]

for i, (title, body) in enumerate(card_data):
    add_rect(s, xs[i], ys[i], 6.1, 2.65,
             fill_rgb=RGBColor(0xEE, 0xF4, 0xFD), line_rgb=MID_BLUE, line_width_pt=1.2)
    # coloured number badge
    add_rect(s, xs[i]+0.1, ys[i]+0.1, 0.45, 0.45, fill_rgb=MID_BLUE)
    add_text(s, str(i+1), xs[i]+0.1, ys[i]+0.07, 0.45, 0.45,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, title, xs[i]+0.65, ys[i]+0.1, 5.3, 0.45,
             font_size=17, bold=True, color=DARK_BLUE)
    add_text(s, body, xs[i]+0.2, ys[i]+0.62, 5.75, 1.9,
             font_size=13.5, color=DARK_TEXT)

footer(s, 3)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — DATASET
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Dataset — openFDA Drug Adverse Events",
             "Publicly accessible, continuously updated safety reports since 1969")

# left block
add_rect(s, 0.35, 1.35, 5.8, 5.5, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Dataset Summary", 0.55, 1.42, 5.4, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Source: FDA MedWatch spontaneous safety report database",
    "Access: openFDA drug/event REST API (no key required)",
    "Format: Nested JSON → Parquet after ingestion",
    "Records processed: 7,573 unique adverse event reports",
    "Fields used: patient age, weight, sex; drug name; reaction MedDRA term; seriousness flags (death, hospitalisation); receipt date",
    "License: Public domain (US government open data)",
], 0.55, 1.9, 5.6, 4.8, font_size=13.5)

# right block — data quality table
add_rect(s, 6.7, 1.35, 6.3, 5.5, fill_rgb=WHITE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Data Quality After Phase 2 Cleaning", 6.9, 1.42, 5.9, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)

rows = [
    ("Field", "Issue", "Rate"),
    ("JSON parse failure", "Bad records routed to bad_records", "~3–5%"),
    ("Patient age", "Null values → median imputed", "~12%"),
    ("Patient weight", "Null values → median imputed", "~28%"),
    ("Duplicate report IDs", "Within single ingest cycle", "<1%"),
]

col_w = [2.3, 2.4, 0.9]
col_x = [6.85, 9.25, 11.75]
row_h = 0.58
row_y0 = 1.95

for ri, row in enumerate(rows):
    fill = MID_BLUE if ri == 0 else (LIGHT_BLUE if ri % 2 == 0 else WHITE)
    txt_col = WHITE if ri == 0 else DARK_TEXT
    for ci, cell in enumerate(row):
        add_rect(s, col_x[ci], row_y0 + ri * row_h, col_w[ci], row_h,
                 fill_rgb=fill, line_rgb=MID_BLUE, line_width_pt=0.5)
        add_text(s, cell, col_x[ci]+0.05, row_y0 + ri * row_h + 0.06,
                 col_w[ci]-0.1, row_h-0.1,
                 font_size=12, bold=(ri == 0), color=txt_col)

footer(s, 4)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "System Architecture",
             "Fully containerised Bronze-Silver-Gold Data Lakehouse")

# pipeline flow diagram (boxes + arrows)
components = [
    ("openFDA\nREST API", 0.5,  2.8, 1.7, 1.4),
    ("Apache\nKafka",     2.7,  2.8, 1.7, 1.4),
    ("HDFS\nBronze",      4.9,  2.8, 1.7, 1.4),
    ("Spark Stream\nSilver", 7.1, 2.8, 1.7, 1.4),
    ("Spark MLlib\nGold",    9.3, 2.8, 1.7, 1.4),
    ("PostgreSQL\n+ Metabase", 11.1, 2.8, 1.9, 1.4),
]

arrow_labels = ["REST", "Kafka\ntopic", "Spark\nstream", "ML\nbatch", "JDBC"]

for idx, (label, lx, ly, lw, lh) in enumerate(components):
    col = MID_BLUE if idx in (1, 2, 3, 4) else DARK_BLUE
    add_rect(s, lx, ly, lw, lh, fill_rgb=col)
    add_text(s, label, lx+0.05, ly+0.18, lw-0.1, lh-0.25,
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# arrows
for i, al in enumerate(arrow_labels):
    ax = components[i][1] + components[i][3] + 0.03
    ay = components[i][2] + components[i][4]/2 - 0.12
    # arrow body
    add_rect(s, ax, ay+0.08, 0.2, 0.08, fill_rgb=ACCENT_GOLD)
    add_text(s, al, ax-0.05, ay-0.28, 0.4, 0.35,
             font_size=9, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# layer labels
add_text(s, "LAYER 1\nIngestion", 0.5, 4.45, 1.7, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "LAYER 2\nBuffering", 2.7, 4.45, 1.7, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "LAYER 3A\nBronze", 4.9, 4.45, 1.7, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "LAYER 3B\nSilver", 7.1, 4.45, 1.7, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "LAYER 4\nGold / ML", 9.3, 4.45, 1.7, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_text(s, "LAYER 5\nReporting", 11.1, 4.45, 1.9, 0.7,
         font_size=10, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# tech justification strip at bottom
add_rect(s, 0.35, 5.35, 12.6, 1.9, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=0.8)
add_text(s, "Technology Choices:", 0.55, 5.42, 3, 0.38,
         font_size=14, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Kafka: log-based replay → exactly-once delivery; decouples API rate limits from Spark speed",
    "HDFS: self-contained, credential-free storage; trivially upgrades to S3/GCS by changing fs.defaultFS",
    "Spark: unified API for Structured Streaming + MLlib; DataFrame API for cleaning; no second ML framework needed",
    "PostgreSQL + Metabase: mature JDBC connector; self-service dashboards; no analyst SQL required",
], 0.55, 5.8, 12.2, 1.4, font_size=11.5)

footer(s, 5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DATA LIFECYCLE & PRE-PROCESSING
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Data Lifecycle & Pre-Processing",
             "From raw openFDA JSON to clean Parquet in HDFS")

# three stage blocks
stages = [
    ("ACQUIRE", "openFDA → Kafka", [
        "Python producer queries drug/event REST endpoint continuously",
        "Each record wrapped in metadata envelope: source_api, ingestion_timestamp, batch_number, raw_data",
        "Published to healthcare_raw Kafka topic (gzip, acks=all)",
        "20 pages × 100 records per ingest cycle = ~2,000 records/run",
    ]),
    ("STORE & CLEAN", "HDFS Bronze → Silver", [
        "Spark Structured Streaming reads Kafka micro-batches (foreachBatch)",
        "Raw envelopes archived to HDFS Bronze (audit-ready Parquet)",
        "Cleaning steps: schema parsing → outlier nullification (negative age/weight) → median imputation → deduplication on report_id",
        "Bad records routed to /healthcare/processed/bad_records",
        "Binary risk_label derived from FDA seriousness flags",
    ]),
    ("PROCESS & VISUALISE", "Gold → PostgreSQL → Metabase", [
        "Spark MLlib batch job reads Silver layer",
        "Feature engineering: StringIndexer → OneHotEncoder → VectorAssembler",
        "Model training: Logistic Regression + Random Forest; best model serialised to HDFS Gold",
        "Predictions + risk summaries written to PostgreSQL via JDBC",
        "Metabase renders per-sex and per-product risk dashboards",
    ]),
]

xs = [0.35, 4.55, 8.75]
for i, (stage, sub, bullets) in enumerate(stages):
    add_rect(s, xs[i], 1.35, 3.8, 5.8, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
    add_rect(s, xs[i], 1.35, 3.8, 0.7, fill_rgb=MID_BLUE)
    add_text(s, stage, xs[i]+0.1, 1.37, 3.6, 0.38,
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, sub, xs[i]+0.1, 1.72, 3.6, 0.35,
             font_size=11, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)
    add_bullet_box(s, bullets, xs[i]+0.15, 2.1, 3.55, 4.9, font_size=12.5)

    if i < 2:
        add_text(s, "▶", xs[i]+3.87, 3.8, 0.5, 0.5,
                 font_size=22, bold=True, color=ACCENT_GOLD, align=PP_ALIGN.CENTER)

footer(s, 6)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ML METHODOLOGY
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "ML Methodology — Risk Classification",
             "Feature engineering + three-model comparison using Spark MLlib")

# Feature engineering table on left
add_rect(s, 0.35, 1.35, 5.8, 2.8, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Feature Engineering Pipeline", 0.55, 1.42, 5.4, 0.45,
         font_size=15, bold=True, color=DARK_BLUE)

feat_rows = [
    ("Type", "Input Columns", "Transformation"),
    ("Numeric", "patient_age, patient_weight, batch_number", "Median imputation → VectorAssembler"),
    ("Categorical", "patient_sex, medicinal_product, reaction, receipt_date", "StringIndexer (handleInvalid=keep) → OneHotEncoder → VectorAssembler"),
]
col_w2 = [1.1, 2.4, 2.1]
col_x2 = [0.4, 1.55, 4.0]
for ri, row in enumerate(feat_rows):
    fill2 = MID_BLUE if ri == 0 else (WHITE if ri % 2 else LIGHT_BLUE)
    tc = WHITE if ri == 0 else DARK_TEXT
    for ci, cell in enumerate(row):
        rh = 0.52 if ri == 0 else 0.72
        add_rect(s, col_x2[ci], 1.9 + sum([0.52 if r == 0 else 0.72 for r in range(ri)]),
                 col_w2[ci], rh, fill_rgb=fill2, line_rgb=MID_BLUE, line_width_pt=0.5)
        add_text(s, cell,
                 col_x2[ci]+0.04, 1.9 + sum([0.52 if r == 0 else 0.72 for r in range(ri)]) + 0.04,
                 col_w2[ci]-0.08, rh-0.08,
                 font_size=10.5, bold=(ri == 0), color=tc)

# Model descriptions on right
model_cards = [
    ("Rule-Based Scorer",
     "Deterministic: high-risk if any FDA seriousness flag = true (serious, death, hospitalisation). Perfectly interpretable, zero training required. Serves as ground-truth baseline."),
    ("Logistic Regression",
     "Parametric probabilistic baseline. Produces calibrated probability estimates for continuous risk ranking. Interpretable via feature coefficients. Efficient on sparse one-hot vectors."),
    ("Random Forest",
     "Primary ML challenger. Captures non-linear drug-reaction-age interactions. Robust to outliers (bootstrap averaging). Handles class imbalance naturally. Best ML model (F1 = 0.86)."),
]

for i, (title, body) in enumerate(model_cards):
    add_rect(s, 6.7, 1.35 + i*1.9, 6.3, 1.75,
             fill_rgb=RGBColor(0xEE, 0xF4, 0xFD), line_rgb=MID_BLUE, line_width_pt=1)
    add_text(s, title, 6.9, 1.38 + i*1.9, 5.9, 0.42,
             font_size=15, bold=True, color=DARK_BLUE)
    add_text(s, body, 6.9, 1.8 + i*1.9, 5.9, 1.2,
             font_size=12.5, color=DARK_TEXT)

# evaluation metric note
add_rect(s, 0.35, 4.3, 12.6, 1.25, fill_rgb=RGBColor(0xFF, 0xF3, 0xCD),
         line_rgb=ACCENT_GOLD, line_width_pt=1.2)
add_text(s, "Primary Metric: Weighted F1 Score",
         0.55, 4.35, 5, 0.42, font_size=14, bold=True, color=DARK_BLUE)
add_text(s,
         "F1 balances precision and recall — critical in healthcare where both false positives (unnecessary clinical alerts) and false negatives (missed serious events) carry real costs. "
         "Evaluation via Spark MulticlassClassificationEvaluator on an 80/20 train-test split.",
         0.55, 4.75, 12.2, 0.75, font_size=12.5, color=DARK_TEXT)

footer(s, 7)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — RESULTS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Results & Performance",
             "7,573 adverse event records — risk classification outcomes")

# Model performance table
add_text(s, "Model Performance (N = 7,573)", 0.35, 1.35, 5.9, 0.45,
         font_size=15, bold=True, color=DARK_BLUE)

perf_rows = [
    ("Model", "Accuracy", "W. Precision", "W. Recall", "W. F1"),
    ("Logistic Regression", "0.84", "0.83", "0.84", "0.83"),
    ("Random Forest",       "0.87", "0.86", "0.87", "0.86 ★"),
    ("Rule-Based Scorer",   "1.00", "1.00", "1.00", "1.00"),
]
pcol_w = [2.0, 1.05, 1.2, 1.1, 0.9]
pcol_x = [0.35, 2.4, 3.5, 4.75, 5.9]

for ri, row in enumerate(perf_rows):
    rh = 0.48
    for ci, cell in enumerate(row):
        fill3 = MID_BLUE if ri == 0 else (RGBColor(0xE8, 0xF5, 0xE9) if ri == 2 else (LIGHT_BLUE if ri % 2 == 0 else WHITE))
        tc = WHITE if ri == 0 else DARK_TEXT
        add_rect(s, pcol_x[ci], 1.82 + ri*rh, pcol_w[ci], rh,
                 fill_rgb=fill3, line_rgb=MID_BLUE, line_width_pt=0.5)
        add_text(s, cell, pcol_x[ci]+0.04, 1.84 + ri*rh,
                 pcol_w[ci]-0.08, rh-0.08,
                 font_size=12, bold=(ri == 0 or (ri == 2 and ci == 4)),
                 color=tc, align=PP_ALIGN.CENTER)

# risk distribution summary
add_rect(s, 0.35, 4.05, 6.2, 2.9, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Risk Distribution", 0.55, 4.1, 5.8, 0.42,
         font_size=15, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "49.1% (3,708 / 7,573) classified as HIGH-RISK (label = 1)",
    "50.9% (3,865 / 7,573) classified as LOW-RISK (label = 0)",
    "Near-equal split reflects MedWatch reporting bias toward serious events",
    "Average risk probability across all records: 0.491",
], 0.55, 4.6, 5.9, 2.3, font_size=13.5)

# per-sex table
add_text(s, "Risk by Patient Sex", 6.9, 1.35, 6.0, 0.42,
         font_size=15, bold=True, color=DARK_BLUE)
sex_rows = [
    ("Sex Group", "Total", "High-Risk", "Low-Risk", "Avg Prob"),
    ("Female (2)", "4,480", "1,915", "2,565", "0.435"),
    ("Male (1)",   "2,326", "1,192", "1,134", "0.511"),
    ("Unknown (0)", "720",  "574",   "146",   "0.767"),
]
sco_w = [1.5, 0.85, 0.95, 0.9, 0.9]
sco_x = [6.9, 8.45, 9.35, 10.35, 11.3]
for ri, row in enumerate(sex_rows):
    rh2 = 0.46
    for ci, cell in enumerate(row):
        fill4 = MID_BLUE if ri == 0 else (LIGHT_BLUE if ri % 2 == 0 else WHITE)
        tc = WHITE if ri == 0 else DARK_TEXT
        add_rect(s, sco_x[ci], 1.82 + ri*rh2, sco_w[ci], rh2,
                 fill_rgb=fill4, line_rgb=MID_BLUE, line_width_pt=0.5)
        add_text(s, cell, sco_x[ci]+0.03, 1.84+ri*rh2,
                 sco_w[ci]-0.06, rh2-0.08,
                 font_size=11.5, bold=(ri==0), color=tc, align=PP_ALIGN.CENTER)

# per-product highlights
add_rect(s, 6.9, 4.05, 6.1, 2.9, fill_rgb=RGBColor(0xEE, 0xF4, 0xFD),
         line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Top Product Risk Highlights", 7.1, 4.1, 5.7, 0.42,
         font_size=15, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Lipitor (atorvastatin): 507 reports, 97.6% high-risk (avg prob 0.929)",
    "Gilenya (fingolimod): 54 reports, 94.4% high-risk (avg prob 0.900)",
    "Pradaxa (dabigatran): 68 reports, 76.5% high-risk (avg prob 0.738)",
    "Warfarin: 15 reports, 80% high-risk — consistent with bleeding profile",
    "Letairis: 1,444 reports, only 6.8% high-risk (low severity reports)",
    "Claritin / Oxytrol: 0% high-risk across all reports",
], 7.1, 4.58, 5.8, 2.3, font_size=12.5)

footer(s, 8)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — COMMUNITY CONTRIBUTION & ETHICS
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Community Contribution & Ethics",
             "Open-source, reproducible, and ethically grounded")

# left col — contributions
add_rect(s, 0.35, 1.35, 6.0, 5.5, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Open-Source Contributions", 0.55, 1.42, 5.6, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Full source code publicly available on GitHub (github.com/2311-ops/bigdata_healthcare_project)",
    "Reusable reference implementation: any openFDA-based analytics project can fork and adapt",
    "Docker Compose setup reproducible in minutes on any machine with Docker installed",
    "Bronze-Silver-Gold HDFS pattern demonstrated outside cloud — useful for institutions without cloud budgets",
    "Metabase dashboard templates exportable for pharma surveillance teams",
    "Educational value: self-contained end-to-end example for CSCI461 and similar courses",
], 0.55, 1.95, 5.8, 4.7, font_size=13.5)

# right col — ethics
add_rect(s, 6.95, 1.35, 6.0, 5.5, fill_rgb=RGBColor(0xFD, 0xF0, 0xE8),
         line_rgb=RGBColor(0xD9, 0x7B, 0x2B), line_width_pt=1)
add_text(s, "Ethical Considerations", 7.15, 1.42, 5.6, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Patient Privacy: openFDA data contains no direct identifiers; combination of rare fields may enable re-identification — platform designed for aggregate population-level analysis only",
    "HIPAA Compliance: any extension to identifiable data requires IRB approval + HIPAA Safe Harbor de-identification + HDFS role-based access control",
    "Clinical Disclaimer: risk labels reflect reporter-perceived seriousness, not validated clinical outcomes; must NOT be used for individual patient decisions",
    "Algorithmic Fairness: unknown-sex records show 79.7% high-risk rate vs 42.7% for females — bias likely from reporting practices, not clinical reality; formal equalised-odds audit recommended",
    "Data Integrity: bad records partitioned and preserved; no silent data loss in the pipeline",
], 7.15, 1.95, 5.8, 4.7, font_size=13.5)

footer(s, 9)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CONCLUSION & FUTURE WORK
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, 13.33, 7.5, fill_rgb=WHITE)
slide_header(s, "Conclusion & Future Work",
             "Summary of achievements and next steps")

# conclusions on left
add_rect(s, 0.35, 1.35, 5.9, 5.5, fill_rgb=LIGHT_BLUE, line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "What We Achieved", 0.55, 1.42, 5.5, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "Ingested 7,573 adverse event records via openFDA → Kafka streaming pipeline",
    "Implemented full Bronze-Silver-Gold HDFS data lakehouse with Spark Structured Streaming",
    "Rule-based risk scorer: perfect classification consistency (F1 = 1.00)",
    "Random Forest: W. F1 = 0.86 — outperforms Logistic Regression (0.83)",
    "49.1% of records classified as high-risk; meaningful variation across drugs and demographics",
    "Fully containerised platform — reproducible on any machine with Docker",
    "Metabase dashboards for self-service pharmacovigilance analytics",
    "Paper addresses all rubric phases: idea, design, data lifecycle, methodology, results, ethics",
], 0.55, 1.95, 5.7, 4.7, font_size=13)

# future work on right
add_rect(s, 6.9, 1.35, 6.1, 5.5, fill_rgb=RGBColor(0xEE, 0xF4, 0xFD),
         line_rgb=MID_BLUE, line_width_pt=1)
add_text(s, "Future Directions", 7.1, 1.42, 5.7, 0.45,
         font_size=16, bold=True, color=DARK_BLUE)
add_bullet_box(s, [
    "NLP feature extraction from narrative adverse event descriptions (Spark NLP / BERT)",
    "Extend pipeline to additional openFDA endpoints: device events, drug recalls",
    "Formal algorithmic fairness audit using equalised-odds metrics across demographic subgroups",
    "Graph Neural Networks for drug-drug interaction signal detection",
    "RAG (Retrieval-Augmented Generation) over the knowledge graph for natural language pharmacovigilance summaries",
    "Production deployment with HDFS replication factor = 3, Kafka partitions = 10+, cloud OLAP for PostgreSQL",
    "Integrate real-time alerting when high-risk thresholds exceeded for specific products",
], 7.1, 1.95, 5.8, 4.7, font_size=13)

footer(s, 10)


# ── save ────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\nowyo\OneDrive - Nile University\Attachments\year 3 sem 2\intro to bigdata\bigdata_healthcare_project\docs\Healthcare_Analytics_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
